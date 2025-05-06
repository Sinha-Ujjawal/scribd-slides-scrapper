#!/usr/bin/python3
import os
import argparse
import logging
from glob import glob
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

class ImagePosn:
    __slots__ = ["img_path", "left", "top", "width", "height"]
    def __init__(self, *, img_path, left, top, width, height):
        self.img_path = img_path
        self.left     = left
        self.top      = top
        self.width    = width
        self.height   = height

def process_image(
    *,
    slide_width,
    slide_height,
    img_path,
    resize_scale,
):
    img_path = os.path.abspath(img_path)
    with Image.open(img_path) as img:
        img = img.resize(tuple(map(lambda x: int(x * resize_scale), img.size)),
                         Image.ANTIALIAS)
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        logging.info(f"Saving scaled down (factor: {resize_scale}) image: {img_path} to temp path: {temp_file.name}")
        img.save(temp_file.name, "PNG")
        temp_img_path = temp_file.name
        img_width_px, img_height_px = img.size
        img_dpi = img.info.get("dpi", (96, 96))
        img_width_in = img_width_px / img_dpi[0]
        img_height_in = img_height_px / img_dpi[1]
    # Convert to pptx units (EMU)
    img_width = Inches(img_width_in)
    img_height = Inches(img_height_in)
    # Scale to fit slide while maintaining aspect ratio
    scale = min(slide_width / img_width, slide_height / img_height)
    final_width = int(img_width * scale)
    final_height = int(img_height * scale)
    # Center image on slide
    left = int((slide_width - final_width) / 2)
    top = int((slide_height - final_height) / 2)
    return ImagePosn(
        img_path = temp_img_path,
        left     = left,
        top      = top,
        width    = final_width,
        height   = final_height,
    )

def convert_images_to_pptx(
    *,
    image_paths,
    out,
    resize_scale,
    max_workers,
):
    slide_no_to_img_posn = {}
    try:
        out = os.path.abspath(out)
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        blank_layout = prs.slide_layouts[6]  # blank layout
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futs = {}
            for slide_no, img_path in enumerate(image_paths, 1):
                fut = executor.submit(
                    process_image,
                    slide_width  = slide_width,
                    slide_height = slide_height,
                    img_path     = img_path,
                    resize_scale = resize_scale,
                )
                futs[fut] = slide_no
            for fut in as_completed(futs):
                slide_no = futs[fut]
                img_posn = fut.result()
                slide_no_to_img_posn[slide_no] = img_posn
        for slide_no in range(1, len(image_paths) + 1):
            img_posn = slide_no_to_img_posn[slide_no]
            logging.info(f"Adding temporary image: {img_posn.img_path} to slide: {slide_no}")
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_posn.img_path,
                img_posn.left,
                img_posn.top,
                width=img_posn.width,
                height=img_posn.height,
            )
            logging.info(f"Added temporary image: {img_posn.img_path} to slide: {slide_no}")

        logging.info(f"Saving presentation to {out}")
        prs.save(out)
        logging.info(f"Saved presentation to {out}")
    finally:
        for img_posn in slide_no_to_img_posn.values():
            temp_file = img_posn.img_path
            logging.info(f"Removing temp file: {temp_file}")
            os.remove(temp_file)
            logging.info(f"Removed temp file: {temp_file}")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--out",
        required=True,
        help="Path of the pptx output",
    )
    parser.add_argument(
        "--max-workers",
        required=False,
        type=int,
        default=4,
        help="Max. no of workers to use for multithreading",
    )
    parser.add_argument(
        "--scale",
        required=False,
        type=float,
        default=1.0,
        help="Scale factor to use for all images",
    )
    parser.add_argument(
        "--image-paths",
        nargs="+",
        required=True,
        help="List of image paths",
    )
    args = parser.parse_args()
    logging.basicConfig(
        format='%(asctime)s - %(levelname)s - %(message)s',
        level=logging.INFO,
        datefmt='%Y-%m-%d %H:%M:%S',
    )
    convert_images_to_pptx(
        image_paths  = args.image_paths,
        out          = args.out,
        max_workers  = args.max_workers,
        resize_scale = args.scale,
    )

if __name__ == "__main__":
    main()
